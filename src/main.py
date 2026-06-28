import os
import re
import math
import requests
from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

def parse_markdown(file_path):
    """
    Markdownファイルを解析して各セクションのテキストや表データを抽出する
    """
    with open(file_path, "r", encoding="utf-8") as f:
        content = f.read()
    
    sections = {}
    
    # 「更新情報無し」が含まれるか、空の場合は早期リターン
    if "更新情報無し" in content or not content.strip():
        sections['no_update'] = True
        return sections
        
    # 結論
    conclusion_match = re.search(r"## a\) 【結論】\n(.*?)(?=\n## b\)|$)", content, re.DOTALL)
    if conclusion_match:
        sections['conclusion'] = conclusion_match.group(1).strip()
        
    # 詳細 (表)
    details_match = re.search(r"## b\) 【詳細】\n(.*?)(?=\n## c\)|$)", content, re.DOTALL)
    if details_match:
        table_text = details_match.group(1).strip()
        lines = [line.strip() for line in table_text.split("\n") if line.strip()]
        
        rows = []
        for line in lines:
            if line.startswith("|") and not re.match(r"^\|[\s:| \- ]+$", line):
                parts = [p.strip() for p in line.split("|")[1:-1]]
                rows.append(parts)
        if len(rows) > 0:
            sections['details_header'] = rows[0]
            sections['details_rows'] = rows[1:]
            
    # 地図的分析
    map_match = re.search(r"## c\) 【地図(?:的)?分析】\n(.*?)(?=\n## d\)|$)", content, re.DOTALL)

    if map_match:
        sections['map_analysis'] = map_match.group(1).strip()
        
    # 影響分析
    impact_match = re.search(r"## d\) 【影響分析】\n(.*?)(?=\n## e\)|$)", content, re.DOTALL)
    if impact_match:
        sections['impact_analysis'] = impact_match.group(1).strip()
        
    # 参照ソース
    sources_match = re.search(r"## e\) 【参照ソース】\n(.*?)$", content, re.DOTALL)
    if sources_match:
        sections['sources'] = sources_match.group(1).strip()
        
    return sections

def deg2num(lat_deg, lon_deg, zoom):
    """
    緯度経度からOSMタイル座標に変換する
    """
    lat_rad = math.radians(lat_deg)
    n = 2.0 ** zoom
    xtile = (lon_deg + 180.0) / 360.0 * n
    ytile = (1.0 - math.asinh(math.tan(lat_rad)) / math.pi) / 2.0 * n
    return xtile, ytile

def get_map_image_from_osm(lat_min, lat_max, lon_min, lon_max, zoom=13):
    """
    指定緯度経度範囲のOpenStreetMapタイル画像を自動ロード・結合して背景地図を生成する
    """
    x_min, y_min = deg2num(lat_max, lon_min, zoom)
    x_max, y_max = deg2num(lat_min, lon_max, zoom)

    x_start = int(math.floor(x_min))
    x_end = int(math.floor(x_max))
    y_start = int(math.floor(y_min))
    y_end = int(math.floor(y_max))

    width = (x_end - x_start + 1) * 256
    height = (y_end - y_start + 1) * 256

    map_img = Image.new('RGB', (width, height))
    headers = {"User-Agent": "IkariCompetitiveMapBot/0.2 (antigravity)"}

    # タイル画像のロードとマージ
    for x in range(x_start, x_end + 1):
        for y in range(y_start, y_end + 1):
            url = f"https://tile.openstreetmap.org/{zoom}/{x}/{y}.png"
            try:
                res = requests.get(url, headers=headers, timeout=5)
                if res.status_code == 200:
                    from io import BytesIO
                    tile = Image.open(BytesIO(res.content))
                else:
                    tile = Image.new('RGB', (256, 256), (245, 245, 245))
            except Exception:
                tile = Image.new('RGB', (256, 256), (245, 245, 245))
                
            px = (x - x_start) * 256
            py = (y - y_start) * 256
            map_img.paste(tile, (px, py))

    # ピクセル座標変換用のインナー関数
    def get_pixel_coords(lat, lon):
        xtile, ytile = deg2num(lat, lon, zoom)
        px = int((xtile - x_start) * 256)
        py = int((ytile - y_start) * 256)
        return px, py

    return map_img, get_pixel_coords

def haversine_distance(lat1, lon1, lat2, lon2):
    """
    緯度経度から2点間のハバース距離(km)を計算する
    """
    R = 6371.0
    phi1 = math.radians(lat1)
    phi2 = math.radians(lat2)
    delta_phi = math.radians(lat2 - lat1)
    delta_lambda = math.radians(lon2 - lon1)
    
    a = math.sin(delta_phi / 2.0)**2 + math.cos(phi1) * math.cos(phi2) * math.sin(delta_lambda / 2.0)**2
    c = 2.0 * math.atan2(math.sqrt(a), math.sqrt(1.0 - a))
    return R * c

def normalize_store_name(name):
    """
    店舗名を比較用に正規化する（スーパー特有の名詞や広域の都市名、店表記、括弧を除去、アルファベット表記揺れの統一）
    """
    # OASIS表記の統一
    name = name.replace("OASIS", "オアシス").replace("oasis", "オアシス")
    
    name = name.split('\n')[0]
    for word in ["スーパー", "ストアー", "ストア", "マーケット", "百貨店", "店", "(仮称)", "（仮称）", "大阪", "兵庫", "京都", "神戸"]:
        name = name.replace(word, "")
    # 空白を除去
    import re
    name = re.sub(r'[\s　]', '', name)
    return name

def filter_stores_for_map(sections, db_data, mode="within_2km"):
    """
    レポートに掲載された競合店と、そのいかりスーパー店舗を抽出する（モードにより半径2km制限を適用）
    """
    details_rows = sections.get('details_rows', [])
    mentioned_names = [row[1] for row in details_rows if len(row) > 1]
    
    ikari_all = {name: tuple(info["coords"]) for name, info in db_data.get("ikari_stores", {}).items()}
    competitors_all = {name: tuple(info["coords"]) for name, info in db_data.get("competitors", {}).items()}
    
    filtered_competitors = {}
    filtered_ikari = {}
    
    if mode == "no_limit":
        # 距離制限なし: レポートに登場した競合店舗すべてと、いかりスーパー全店舗を表示
        for comp_db_name, comp_coords in competitors_all.items():
            norm_db = normalize_store_name(comp_db_name)
            is_mentioned = False
            for name in mentioned_names:
                norm_mentioned = normalize_store_name(name)
                is_matched = False
                min_len = 2
                for i in range(len(norm_db) - min_len + 1):
                    sub = norm_db[i:i+min_len]
                    if sub in norm_mentioned:
                        is_matched = True
                        break
                if is_matched:
                    is_mentioned = True
                    break
            if is_mentioned:
                filtered_competitors[comp_db_name] = comp_coords
        return ikari_all, filtered_competitors
        
    for comp_db_name, comp_coords in competitors_all.items():
        norm_db = normalize_store_name(comp_db_name)
        is_mentioned = False
        for name in mentioned_names:
            norm_mentioned = normalize_store_name(name)
            # 2文字以上の共通部分文字列があるか判定（部分一致の条件を緩和）
            is_matched = False
            min_len = 2
            for i in range(len(norm_db) - min_len + 1):
                sub = norm_db[i:i+min_len]
                if sub in norm_mentioned:
                    is_matched = True
                    break
            if is_matched:
                is_mentioned = True
                break
                
        if is_mentioned:
            # 半径2km以内にいかりスーパーがあるか判定
            nearby_ikari = {}
            for ikari_name, ikari_coords in ikari_all.items():
                dist = haversine_distance(comp_coords[0], comp_coords[1], ikari_coords[0], ikari_coords[1])
                if dist <= 2.0:
                    nearby_ikari[ikari_name] = ikari_coords
            
            # いかりスーパーが半径2km以内にある競合、およびそのいかり店舗のみを対象とする
            if len(nearby_ikari) > 0:
                filtered_competitors[comp_db_name] = comp_coords
                for k, v in nearby_ikari.items():
                    filtered_ikari[k] = v
                    
    # 該当店舗がない（または更新情報無しなどの）場合は、デフォルトでいかりスーパー全店を表示
    if not filtered_competitors:
        filtered_ikari = ikari_all
        
    return filtered_ikari, filtered_competitors


def generate_visual_map(output_path, ikari_stores, competitors):
    """
    OpenStreetMapを背景にマージした正確な店舗位置関係画像を自動生成する
    """
    coords = []
    for latlon in ikari_stores.values():
        coords.append(latlon)
    for latlon in competitors.values():
        coords.append(latlon)
        
    if len(coords) > 0:
        lats = [c[0] for c in coords]
        lons = [c[1] for c in coords]
        # スパンに応じてマージンを動的に調整
        lat_span = max(lats) - min(lats)
        lon_span = max(lons) - min(lons)
        max_span = max(lat_span, lon_span)
        
        # スパンに応じてズームとマージンを自動調整
        if max_span > 0.4:
            zoom = 10
            lat_margin, lon_margin = 0.05, 0.07
        elif max_span > 0.15:
            zoom = 11
            lat_margin, lon_margin = 0.03, 0.04
        elif max_span > 0.06:
            zoom = 12
            lat_margin, lon_margin = 0.02, 0.025
        else:
            zoom = 13
            lat_margin, lon_margin = 0.015, 0.020
            
        lat_min, lat_max = min(lats) - lat_margin, max(lats) + lat_margin
        lon_min, lon_max = min(lons) - lon_margin, max(lons) + lon_margin
    else:
        # デフォルトの阪神間エリア範囲
        lat_min, lat_max = 34.710, 34.815
        lon_min, lon_max = 135.315, 135.425
        zoom = 13
    
    # 1. 地図画像ロード (動的に算出した zoom を渡す)
    map_img, get_pixel_coords = get_map_image_from_osm(lat_min, lat_max, lon_min, lon_max, zoom=zoom)
    
    # 2. マーカー描画
    draw = ImageDraw.Draw(map_img)
    
    # Windows標準フォントの設定 (文字化け対策)
    try:
        font = ImageFont.truetype("msgothic.ttc", 12)
        font_bold = ImageFont.truetype("msgothic.ttc", 13)
    except Exception:
        font = ImageFont.load_default()
        font_bold = ImageFont.load_default()
        
    # いかりスーパーを描画 (ブルー、ゴールド枠)
    for name, (lat, lon) in ikari_stores.items():
        px, py = get_pixel_coords(lat, lon)
        r = 10
        # ゴールド外枠
        draw.ellipse([px-r-2, py-r-2, px+r+2, py+r+2], fill=(212, 175, 55))
        # ブルー丸
        draw.ellipse([px-r, py-r, px+r, py+r], fill=(0, 91, 172))
        # テキストラベル背景
        text_w = draw.textlength(name, font=font_bold)
        draw.rectangle([px - text_w/2 - 4, py - r - 20, px + text_w/2 + 4, py - r - 3], fill=(255, 255, 255), outline=(0, 91, 172), width=1)
        # ラベル描画
        draw.text((px, py - r - 19), name, fill=(0, 91, 172), font=font_bold, anchor="ma")

    # 競合店舗を描画 (赤: ヤマダ、オレンジ: その他)
    for name, (lat, lon) in competitors.items():
        px, py = get_pixel_coords(lat, lon)
        r = 8
        color = (255, 0, 0) if "ヤマダ" in name else (255, 140, 0)
        # 黒外枠
        draw.ellipse([px-r-1, py-r-1, px+r+1, py+r+1], fill=(0, 0, 0))
        # 塗りつぶし
        draw.ellipse([px-r, py-r, px+r, py+r], fill=color)
        
        # ラベル名
        lines = name.split("\n")
        text_y = py + r + 2
        for line in lines:
            text_w = draw.textlength(line, font=font)
            draw.rectangle([px - text_w/2 - 3, text_y, px + text_w/2 + 3, text_y + 13], fill=(255, 255, 255), outline=(100, 100, 100), width=1)
            draw.text((px, text_y), line, fill=(0, 0, 0), font=font, anchor="ma")
            text_y += 14

    # 凡例（Legend）を描画 (左上に枠付きで描画)
    draw.rectangle([10, 10, 260, 100], fill=(255, 255, 255, 230), outline=(150, 150, 150), width=1)
    # いかり
    draw.ellipse([20-6, 25-6, 20+6, 25+6], fill=(0, 91, 172), outline=(212, 175, 55), width=1)
    draw.text((35, 18), "いかりスーパー", fill=(0, 91, 172), font=font_bold)
    # ヤマダ
    draw.ellipse([20-6, 50-6, 20+6, 50+6], fill=(255, 0, 0), outline=(0, 0, 0), width=1)
    draw.text((35, 43), "ヤマダストアー (Tier 1)", fill=(0, 0, 0), font=font)
    # その他
    draw.ellipse([20-6, 75-6, 20+6, 75+6], fill=(255, 140, 0), outline=(0, 0, 0), width=1)
    draw.text((35, 68), "その他重点競合 (OK/オアシス)", fill=(0, 0, 0), font=font)

    # 3. 保存
    os.makedirs(os.path.dirname(output_path), exist_ok=True)
    map_img.save(output_path, "PNG")
    print(f"OSMマージ地図画像を生成・保存しました: {output_path}")

def generate_dynamic_web_page(output_path, db_data, md_path, template_path):
    """
    マージされた店舗データとレポートMarkdownのデータを
    dashboard_template.html 内に埋め込み、統合されたWebダッシュボードを生成する
    """
    import json
    
    if not os.path.exists(template_path):
        print(f"Error: Template file {template_path} not found.")
        return
        
    if not os.path.exists(md_path):
        print(f"Error: Markdown file {md_path} not found.")
        return

    # 1. データベースデータのJSON文字列化
    db_content = json.dumps(db_data, ensure_ascii=False, indent=2)
        
    # 2. レポートMarkdownをパース
    sections = parse_markdown(md_path)
    report_json = json.dumps(sections, ensure_ascii=False)
    
    # 3. テンプレートの読み込み
    with open(template_path, "r", encoding="utf-8") as f:
        html_content = f.read()
        
    # 4. プレースホルダを置換
    dynamic_html = html_content.replace("/*STORE_DATA_PLACEHOLDER*/", db_content.strip())
    dynamic_html = dynamic_html.replace("/*REPORT_DATA_PLACEHOLDER*/", report_json.strip())
    
    # 5. HTMLの保存
    os.makedirs(os.path.dirname(output_path), exist_ok=True)
    with open(output_path, "w", encoding="utf-8") as f:
        f.write(dynamic_html)
        
    print(f"統合ダッシュボードWebページを生成・保存しました: {output_path}")

def create_presentation(sections, output_path, map_image_path):
    """
    抽出したセクションデータからPowerPoint(16:9)を生成する
    """
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)
    
    THEME_COLOR = RGBColor(0, 91, 172)
    GOLD = RGBColor(212, 175, 55)
    WHITE = RGBColor(255, 255, 255)
    BLACK = RGBColor(30, 30, 30)
    GRAY = RGBColor(245, 245, 245)
    
    # 1. タイトルスライド
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = THEME_COLOR
    
    txBox = slide.shapes.add_textbox(Inches(1), Inches(2.2), Inches(11.33), Inches(2.5))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "いかりスーパー周辺 競合環境調査レポート"
    p.font.bold = True
    p.font.size = Pt(40)
    p.font.color.rgb = GOLD
    p.alignment = PP_ALIGN.CENTER
    
    p2 = tf.add_paragraph()
    p2.text = "2026年6月度 競合調査結果報告書"
    p2.font.size = Pt(22)
    p2.font.color.rgb = WHITE
    p2.alignment = PP_ALIGN.CENTER
    p2.space_before = Pt(20)
    
    # 「更新情報無し」の場合の処理
    if sections.get('no_update'):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        title_box = slide.shapes.add_textbox(Inches(0.75), Inches(0.5), Inches(11.83), Inches(0.8))
        title_box.text_frame.word_wrap = True
        p = title_box.text_frame.paragraphs[0]
        p.text = "周辺環境の更新情報について"
        p.font.bold = True
        p.font.size = Pt(28)
        p.font.color.rgb = THEME_COLOR
        
        content_box = slide.shapes.add_textbox(Inches(0.75), Inches(2.5), Inches(11.83), Inches(3.0))
        tf_content = content_box.text_frame
        tf_content.word_wrap = True
        p_item = tf_content.paragraphs[0]
        p_item.text = "本日の調査において、前日からの周辺環境（競合の新規出店・改装計画・その他動向）の更新情報はありませんでした。"
        p_item.font.size = Pt(20)
        p_item.font.color.rgb = BLACK
        p_item.alignment = PP_ALIGN.CENTER
        
        prs.save(output_path)
        print(f"Presentation saved (No Updates) to: {output_path}")
        return
        
    # 2. 結論スライド
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    title_box = slide.shapes.add_textbox(Inches(0.75), Inches(0.5), Inches(11.83), Inches(0.8))
    title_box.text_frame.word_wrap = True
    p = title_box.text_frame.paragraphs[0]
    p.text = "a) 【結論】"
    p.font.bold = True
    p.font.size = Pt(28)
    p.font.color.rgb = THEME_COLOR
    
    content_box = slide.shapes.add_textbox(Inches(0.75), Inches(1.5), Inches(11.83), Inches(5.0))
    tf_content = content_box.text_frame
    tf_content.word_wrap = True
    
    conclusion_text = sections.get('conclusion', '')
    conclusion_lines = [line.strip().replace('* ', '').replace('- ', '') for line in conclusion_text.split("\n") if line.strip()]
    
    for idx, item in enumerate(conclusion_lines):
        p_item = tf_content.add_paragraph() if idx > 0 else tf_content.paragraphs[0]
        p_item.text = "・ " + item
        p_item.font.size = Pt(18)
        p_item.font.color.rgb = BLACK
        p_item.space_after = Pt(16)
        
    # 3. 詳細スライド（表データ、6行ずつ複数スライドに分割表示）
    if 'details_rows' in sections:
        rows_data = sections['details_rows']
        headers = sections['details_header']
        rows_per_slide = 6
        
        for slide_idx in range((len(rows_data) + rows_per_slide - 1) // rows_per_slide):
            slide_rows_data = rows_data[slide_idx * rows_per_slide : (slide_idx + 1) * rows_per_slide]
            
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            
            title_box = slide.shapes.add_textbox(Inches(0.75), Inches(0.5), Inches(11.83), Inches(0.8))
            title_box.text_frame.word_wrap = True
            p = title_box.text_frame.paragraphs[0]
            p.text = f"b) 【詳細】 ({slide_idx + 1}ページ目)"
            p.font.bold = True
            p.font.size = Pt(28)
            p.font.color.rgb = THEME_COLOR
            
            rows = len(slide_rows_data) + 1
            cols = len(headers)
            
            left = Inches(0.75)
            top = Inches(1.4)
            width = Inches(11.83)
            height = Inches(5.2)
            
            table_shape = slide.shapes.add_table(rows, cols, left, top, width, height)
            table = table_shape.table
            
            col_widths = [Inches(1.0), Inches(1.8), Inches(1.0), Inches(1.0), Inches(1.5), Inches(1.0), Inches(4.53)]
            for c_idx, w in enumerate(col_widths):
                if c_idx < cols:
                    table.columns[c_idx].width = w
                
            for col_idx, header in enumerate(headers):
                cell = table.cell(0, col_idx)
                cell.text = header
                cell.fill.solid()
                cell.fill.fore_color.rgb = THEME_COLOR
                for paragraph in cell.text_frame.paragraphs:
                    paragraph.font.bold = True
                    paragraph.font.size = Pt(11)
                    paragraph.font.color.rgb = WHITE
                    paragraph.alignment = PP_ALIGN.CENTER
                    
            for row_idx, row_data in enumerate(slide_rows_data):
                for col_idx, cell_value in enumerate(row_data):
                    if col_idx < len(headers):
                        cell = table.cell(row_idx + 1, col_idx)
                        cell.text = cell_value
                        if row_idx % 2 == 0:
                            cell.fill.solid()
                            cell.fill.fore_color.rgb = GRAY
                        for paragraph in cell.text_frame.paragraphs:
                            paragraph.font.size = Pt(9.5)
                            paragraph.font.color.rgb = BLACK
                            
    # 4. 地図的分析スライド
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    title_box = slide.shapes.add_textbox(Inches(0.75), Inches(0.5), Inches(11.83), Inches(0.8))
    title_box.text_frame.word_wrap = True
    p = title_box.text_frame.paragraphs[0]
    p.text = "c) 【地図的分析】"
    p.font.bold = True
    p.font.size = Pt(28)
    p.font.color.rgb = THEME_COLOR
    
    content_box = slide.shapes.add_textbox(Inches(0.75), Inches(1.4), Inches(11.83), Inches(5.3))
    tf_content = content_box.text_frame
    tf_content.word_wrap = True
    
    map_text = sections.get('map_analysis', '')
    map_lines = []
    current_item = ""
    for line in map_text.split("\n"):
        line = line.strip()
        if not line:
            continue
        if re.match(r"^\d+\.", line) or line.startswith("-") or line.startswith("*"):
            if current_item:
                map_lines.append(current_item)
            current_item = line
        else:
            if current_item:
                current_item += "\n" + line
            else:
                current_item = line
    if current_item:
        map_lines.append(current_item)
        
    for idx, item in enumerate(map_lines):
        p_item = tf_content.add_paragraph() if idx > 0 else tf_content.paragraphs[0]
        p_item.text = item
        p_item.font.size = Pt(14)
        p_item.font.color.rgb = BLACK
        p_item.space_after = Pt(14)
        
    # 5. 影響分析スライド
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    title_box = slide.shapes.add_textbox(Inches(0.75), Inches(0.5), Inches(11.83), Inches(0.8))
    title_box.text_frame.word_wrap = True
    p = title_box.text_frame.paragraphs[0]
    p.text = "d) 【影響分析】"
    p.font.bold = True
    p.font.size = Pt(28)
    p.font.color.rgb = THEME_COLOR
    
    content_box = slide.shapes.add_textbox(Inches(0.75), Inches(1.4), Inches(11.83), Inches(5.3))
    tf_content = content_box.text_frame
    tf_content.word_wrap = True
    
    impact_text = sections.get('impact_analysis', '')
    impact_lines = []
    current_item = ""
    for line in impact_text.split("\n"):
        line = line.strip()
        if not line:
            continue
        if re.match(r"^\d+\.", line) or line.startswith("-") or line.startswith("*"):
            if current_item:
                impact_lines.append(current_item)
            current_item = line
        else:
            if current_item:
                current_item += "\n" + line
            else:
                current_item = line
    if current_item:
        impact_lines.append(current_item)
        
    for idx, item in enumerate(impact_lines):
        p_item = tf_content.add_paragraph() if idx > 0 else tf_content.paragraphs[0]
        p_item.text = item
        p_item.font.size = Pt(13.5)
        p_item.font.color.rgb = BLACK
        p_item.space_after = Pt(12)
        
    # 6. 参照ソーススライド
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    title_box = slide.shapes.add_textbox(Inches(0.75), Inches(0.5), Inches(11.83), Inches(0.8))
    title_box.text_frame.word_wrap = True
    p = title_box.text_frame.paragraphs[0]
    p.text = "e) 【参照ソース】"
    p.font.bold = True
    p.font.size = Pt(28)
    p.font.color.rgb = THEME_COLOR
    
    content_box = slide.shapes.add_textbox(Inches(0.75), Inches(1.4), Inches(11.83), Inches(5.3))
    tf_content = content_box.text_frame
    tf_content.word_wrap = True
    
    sources_text = sections.get('sources', '')
    sources_lines = [line.strip().replace('* ', '').replace('- ', '') for line in sources_text.split("\n") if line.strip()]
    
    for idx, item in enumerate(sources_lines):
        p_item = tf_content.add_paragraph() if idx > 0 else tf_content.paragraphs[0]
        p_item.space_after = Pt(8)
        
        # '・ ' を先頭に追加
        run_bullet = p_item.add_run()
        run_bullet.text = "・ "
        run_bullet.font.size = Pt(12)
        run_bullet.font.color.rgb = BLACK
        
        match = re.search(r"\[([^\]]+)\]\((https?://[^\)]+)\)", item)
        if match:
            start, end = match.span()
            link_text = match.group(1)
            url = match.group(2)
            
            # リンクの表記テキスト自体がURLで、その手前にテキストが存在する場合、
            # 手前のテキストをリンクタイトルとして流用する
            if (link_text.startswith("http://") or link_text.startswith("https://")) and start > 0:
                title = item[:start].strip().rstrip(".").rstrip("。").strip()
                if not title:
                    title = link_text
                before = ""
            else:
                title = link_text
                before = item[:start]
                
            after = item[end:]
            
            if before:
                run_before = p_item.add_run()
                run_before.text = before
                run_before.font.size = Pt(12)
                run_before.font.color.rgb = BLACK
                
            run_link = p_item.add_run()
            run_link.text = title
            run_link.font.size = Pt(12)
            run_link.font.color.rgb = RGBColor(0, 102, 204)  # 綺麗な青
            run_link.font.underline = True
            run_link.hyperlink.address = url
            
            if after:
                run_after = p_item.add_run()
                run_after.text = after
                run_after.font.size = Pt(12)
                run_after.font.color.rgb = BLACK
        else:
            # マッチしない場合は通常のテキストとして追加
            run_text = p_item.add_run()
            run_text.text = item
            run_text.font.size = Pt(12)
            run_text.font.color.rgb = BLACK
        
    # 7. ビジュアルマップスライド (OSM地図マージ画像を埋め込み)
    if os.path.exists(map_image_path):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        title_box = slide.shapes.add_textbox(Inches(0.75), Inches(0.5), Inches(11.83), Inches(0.8))
        title_box.text_frame.word_wrap = True
        p = title_box.text_frame.paragraphs[0]
        p.text = "f) 【ビジュアルマップ】"
        p.font.bold = True
        p.font.size = Pt(28)
        p.font.color.rgb = THEME_COLOR
        
        # 元画像のアスペクト比を維持してスライド中央に配置
        from PIL import Image
        with Image.open(map_image_path) as img:
            img_w, img_h = img.size
            
        max_w = Inches(9.5)
        max_h = Inches(5.7)
        
        # 比率を計算してフィットするサイズを算出
        ratio = min(max_w / img_w, max_h / img_h)
        fit_w = img_w * ratio
        fit_h = img_h * ratio
        
        # 表示枠（横9.5, 縦5.7, 左位置1.91, 上位置1.3）の中央に寄せる
        left = Inches(1.91) + (max_w - fit_w) / 2
        top = Inches(1.3) + (max_h - fit_h) / 2
        
        slide.shapes.add_picture(map_image_path, left, top, width=fit_w, height=fit_h)
        
    prs.save(output_path)
    print(f"Presentation saved successfully to: {output_path}")

def parse_ikari_config_md(config_path):
    if not os.path.exists(config_path):
        return {}
        
    ikari_stores = {}
    with open(config_path, "r", encoding="utf-8") as f:
        lines = f.readlines()
        
    table_started = False
    for line in lines:
        line = line.strip()
        if line.startswith("|") and line.endswith("|"):
            if "店舗名" in line and "住所" in line:
                table_started = True
                continue
            if table_started:
                # 区切り行をスキップ
                if re.match(r"^\|[\s:|#-]+\|$", line):
                    continue
                
                cols = [c.strip() for c in line.split("|")][1:-1]
                if len(cols) >= 9:
                    name = cols[0]
                    coords_str = cols[1]
                    if not name:
                        continue
                    
                    # 座標のパース
                    coords = None
                    if coords_str:
                        try:
                            lat_lon = [float(x.strip()) for x in coords_str.split(",")]
                            if len(lat_lon) == 2:
                                coords = tuple(lat_lon)
                        except ValueError:
                            pass
                            
                    ikari_stores[name] = {
                        "coords": coords,
                        "address": cols[2],
                        "phone": cols[3],
                        "parking": cols[4],
                        "hours": cols[5],
                        "features": cols[6],
                        "manager": cols[7] if len(cols) > 7 else "",
                        "image_url": cols[8] if len(cols) > 8 else "",
                        "detail_url": cols[9] if len(cols) > 9 else ""
                    }
    return ikari_stores

def main():
    import json
    import argparse
    from datetime import datetime
    
    parser = argparse.ArgumentParser(description="Generate maps and presentation based on competitive report")
    parser.add_argument("--mode", choices=["within_2km", "no_limit"], default="within_2km", help="Generation mode")
    args = parser.parse_args()
    mode = args.mode
    
    today_str = datetime.today().strftime("%Y_%m_%d")
    md_path = "report/competitive_report_within_2km.md" if mode == "within_2km" else "report/competitive_report_no_limit.md"
    pptx_path = f"report/{today_str}_competitive_report_{mode}.pptx"
    map_image_path = f"report/{today_str}_competitive_map_{mode}.png"
    map_html_path = f"report/{today_str}_competitive_map_{mode}.html"
    db_path = "stores_db.json"
    ikari_config_path = "../ikari_map/config.md"
    
    if not os.path.exists(md_path):
        print(f"Error: {md_path} not found. Cannot convert to PPTX.")
        return
        
    if not os.path.exists(db_path):
        print(f"Error: {db_path} not found. Cannot load store locations.")
        return

    # 1. 基準となる位置情報マスタJSONからデータをロード
    with open(db_path, "r", encoding="utf-8") as f:
        db_data = json.load(f)

    # 2. ikari_map/config.md から最新のいかりスーパー店舗情報を読み込んでマージ
    config_stores = parse_ikari_config_md(ikari_config_path)
    if config_stores:
        print(f"Loaded {len(config_stores)} stores from {ikari_config_path}")
        
        # 既存の ikari_stores 辞書をベースに、config.md の情報で更新
        merged_ikari = {}
        
        # 元の stores_db.json にある店舗を config.md で上書きマージ
        for name, orig_info in db_data.get("ikari_stores", {}).items():
            if name in config_stores:
                cfg_info = config_stores[name]
                # 座標が config.md に定義されている場合はそれを使用し、無ければ元の座標を使う
                coords = list(cfg_info["coords"]) if cfg_info["coords"] else orig_info.get("coords")
                merged_ikari[name] = {
                    "coords": coords,
                    "address": cfg_info["address"] if cfg_info["address"] else orig_info.get("address"),
                    "phone": cfg_info["phone"] if cfg_info["phone"] else orig_info.get("phone"),
                    "parking": cfg_info["parking"] if cfg_info["parking"] else orig_info.get("parking"),
                    "hours": cfg_info["hours"] if cfg_info["hours"] else orig_info.get("hours"),
                    "features": cfg_info["features"] if cfg_info["features"] else orig_info.get("features"),
                    "manager": cfg_info.get("manager") if cfg_info.get("manager") else orig_info.get("manager", ""),
                    "image_url": cfg_info["image_url"] if cfg_info["image_url"] else orig_info.get("image_url"),
                    "detail_url": cfg_info.get("detail_url") if cfg_info.get("detail_url") else orig_info.get("detail_url", ""),
                    "area": orig_info.get("area", "N/A")
                }
            else:
                # config.md に記載がない店舗はそのまま残す
                merged_ikari[name] = orig_info
                
        # config.md にのみ存在する新規店舗があれば追加
        for name, cfg_info in config_stores.items():
            if name not in merged_ikari:
                # 座標がある場合のみ追加
                if cfg_info["coords"]:
                    merged_ikari[name] = {
                        "coords": list(cfg_info["coords"]),
                        "address": cfg_info["address"],
                        "phone": cfg_info["phone"],
                        "parking": cfg_info["parking"],
                        "hours": cfg_info["hours"],
                        "features": cfg_info["features"],
                        "manager": cfg_info.get("manager", ""),
                        "image_url": cfg_info["image_url"],
                        "detail_url": cfg_info.get("detail_url", ""),
                        "area": "N/A"
                    }
                    
        db_data["ikari_stores"] = merged_ikari

    # リスト形式の座標をタプルに変換してマッピングに適用
    ikari_stores = {}
    for name, info in db_data.get("ikari_stores", {}).items():
        if info.get("coords"):
            ikari_stores[name] = tuple(info["coords"])
            
    competitors = {name: tuple(info["coords"]) for name, info in db_data.get("competitors", {}).items()}
    
    # 0. Markdown解析
    sections = parse_markdown(md_path)

    # 1. OpenStreetMapマージによる高精度な背景地図画像の自動生成
    # (更新情報がある場合のみ、スライド用のビジュアルマップ画像を生成)
    if not sections.get('no_update'):
        filtered_ikari, filtered_comp = filter_stores_for_map(sections, db_data, mode=mode)
        generate_visual_map(map_image_path, filtered_ikari, filtered_comp)
    
    # 2. レポート統合型HTMLダッシュボードの自動生成
    generate_dynamic_web_page(map_html_path, db_data, md_path, "dashboard_template.html")
    
    # 3. PowerPoint生成（地図画像をスライドに挿入）
    os.makedirs(os.path.dirname(pptx_path), exist_ok=True)
    create_presentation(sections, pptx_path, map_image_path)
    
    # PPTX生成完了後、スライド埋め込み用の一時マップ画像を削除
    if os.path.exists(map_image_path):
        try:
            os.remove(map_image_path)
            print(f"Temporary map image deleted: {map_image_path}")
        except OSError:
            pass

    # レポート生成の最終成功日を記録
    try:
        from datetime import datetime
        last_success_path = os.path.join(os.path.dirname(pptx_path), "last_success.txt")
        with open(last_success_path, "w", encoding="utf-8") as f:
            f.write(datetime.now().strftime("%Y-%m-%d"))
        print(f"Last success date written to: {last_success_path}")
    except Exception as e:
        print(f"Error writing last success file: {e}")


if __name__ == "__main__":
    main()
